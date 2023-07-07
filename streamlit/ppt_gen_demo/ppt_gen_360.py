from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches
import streamlit as st 
import requests
import json
from datetime import datetime
from PIL import Image


st.set_page_config(page_title="InSource AI PowerPoint Generator", page_icon=":frame_with_picture:")

hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """
st.markdown(hide_streamlit_style, unsafe_allow_html=True) 

def ppt_gen(pptstyle, product, ideas):
    # Open the PowerPoint file
    ppttemplate = "ProductTemplate-" + pptstyle + ".pptx"
    prs = Presentation(ppttemplate)

    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if shape.has_text_frame:
                # Loop through each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Loop through each run in the paragraph
                    for run in paragraph.runs:
                        # Print the text of the run
                        print(run.text)

    # Get slide 1
    slide = prs.slides[0]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Product' in paragraph.text:
                    paragraph.text = product
                if '###Date' in paragraph.text:
                    now = datetime.now()
                    paragraph.text = now.strftime("%Y-%m-%d")

    # Get slide 2
    slide = prs.slides[1]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Slide-1-Subject' in paragraph.text:
                    paragraph.text = ideas["Slide1Subject"]
                if '###Slide-1-Description' in paragraph.text:
                    response_text = f"""
                    You are a treatment sales. You need to sell {product} to your customer.
                    You are writing content on a PowerPoint slide with the subject: {ideas["Slide1Subject"]}.
                    Give the content within 30 words.
                    """
                    url = "http://insource-test-015.southeastasia.cloudapp.azure.com:8000/run"
                    payload = json.dumps({
                    "id": "1",
                    "text": response_text
                    })
                    headers = {
                    'Content-Type': 'application/json'
                    }
                    response = requests.request("POST", url, headers=headers, data=payload)          
                    print(response.text)
                    paragraph.text = json.loads(response.text)["result"]

    # Get slide 3
    slide = prs.slides[2]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Slide-2-Subject' in paragraph.text:
                    paragraph.text = ideas["Slide2Subject"]
                if '###Slide-2-Description' in paragraph.text:
                    response_text = f"""
                    You are a treatment sales. You need to sell {product} to your customer.
                    You are writing content on a PowerPoint slide with the subject: {ideas["Slide2Subject"]}.
                    Give the content within 30 words.
                    """
                    url = "http://insource-test-015.southeastasia.cloudapp.azure.com:8000/run"
                    payload = json.dumps({
                    "id": "1",
                    "text": response_text
                    })
                    headers = {
                    'Content-Type': 'application/json'
                    }
                    response = requests.request("POST", url, headers=headers, data=payload)          
                    print(response.text)
                    paragraph.text = json.loads(response.text)["result"]

    # Get slide 4
    slide = prs.slides[3]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Slide-3-Subject' in paragraph.text:
                    paragraph.text = "Compare Hifu with another treatment"
                if '###Slide-3-Description' in paragraph.text:
                    response_text = f"""
                    You are a treatment sales. You need to sell {product} to your customer.
                    You are writing content on a PowerPoint slide with the subject: {ideas["Slide2Subject"]}.
                    Give the content within 30 words.
                    """
                    url = "http://insource-test-015.southeastasia.cloudapp.azure.com:8000/run"
                    payload = json.dumps({
                    "id": "1",
                    "text": response_text
                    })
                    headers = {
                    'Content-Type': 'application/json'
                    }
                    response = requests.request("POST", url, headers=headers, data=payload)          
                    print(response.text)
                    paragraph.text = json.loads(response.text)["result"]

    # Get slide 5
    slide = prs.slides[4]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Slide-4-Subject' in paragraph.text:
                    paragraph.text = ideas["Slide4Subject"]
                if '###Slide-4-Description' in paragraph.text:
                    response_text = f"""
                    You are a treatment sales. You need to sell {product} to your customer.
                    You are writing content on a PowerPoint slide with the subject: {ideas["Slide4Subject"]}.
                    Give the content within 30 words.
                    """
                    url = "http://insource-test-015.southeastasia.cloudapp.azure.com:8000/run"
                    payload = json.dumps({
                    "id": "1",
                    "text": response_text
                    })
                    headers = {
                    'Content-Type': 'application/json'
                    }
                    response = requests.request("POST", url, headers=headers, data=payload)          
                    print(response.text)
                    paragraph.text = json.loads(response.text)["result"]                


    # Get slide 6
    slide = prs.slides[5]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###Advantage1' in paragraph.text:
                    paragraph.text = ideas["ProductAdvantage1"]
                if '###Advantage2' in paragraph.text:
                    paragraph.text = ideas["ProductAdvantage2"]    
                if '###Advantage3' in paragraph.text:
                    paragraph.text = ideas["ProductAdvantage3"]
                if '###Advantage4' in paragraph.text:
                    paragraph.text = ideas["ProductAdvantage4"]

    # Get slide 8
    slide = prs.slides[7]

    # Find and replace text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if '###ProductSlideSummary' in paragraph.text:
                    paragraph.text = ideas["ProductSlideSummary"]               

    # Get slide 9
    slide = prs.slides[8]

    for shape in slide.shapes:
        if shape.has_chart:
            slide.shapes.remove(shape)

    chart_data = ChartData()
    chart_data.categories = ['A', 'B', 'C', 'D', 'E']
    chart_data.add_series('Test', (0.135, 0.324, 0.180, 0.235, 0.126))

    x, y, cx, cy = Inches(4), Inches(2), Inches(6), Inches(4.5)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Save the presentation
    pptfilename = "ProductDemo__" + product + now.strftime("%Y-%m-%d-%h-%m-%s") + ".pptx"
    prs.save(pptfilename)

    # Add a download button for the user to download the PowerPoint file
    with open(pptfilename, "rb") as f:
        bytes_data = f.read()
        st.download_button(
            label="Download PowerPoint",
            data=bytes_data,
            file_name=pptfilename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

def insource_idea(product):
    response_text = f"""
    You are a treatment sales. You need to design a powerpoint for selling {product} to your customer.
    The powerpoint is having 4 slides. Please describe each of the slides you will talk about.
    The final answer must be in JSON string like:

     "Slide1Subject": <Slide 1 Subject>, "Slide1Description": <Slide 1 Description>, 
     "Slide2Subject": <Slide 2 Subject>, "Slide2Description": <Slide 2 Description>,
     "Slide3Subject": <Slide 3 Subject>, "Slide3Description": <Slide 3 Description>,
     "Slide4Subject": <Slide 4 Subject>, "Slide4Description": <Slide 4 Description>,
     "ProductAdvantage1": <Product Advantage 1>,
     "ProductAdvantage2": <Product Advantage 2>,
     "ProductAdvantage3": <Product Advantage 3>,
     "ProductAdvantage4": <Product Advantage 4>,
     "ProductSlideSummary": <Product Summary in 70 words>
    """
    
    url = "http://insource-test-015.southeastasia.cloudapp.azure.com:8000/run"  
    payload = {  
      "id": "2",  
      "text": response_text  
    }
    response = requests.post(url, json=payload)  
    return response.json()


def main():  
    st.title("InSource AI PowerPoint Generator") 

    col1, col2 = st.columns([1, 3])

    with col1:
        pptstyle = st.selectbox(
            "Style",
            ('ClockBlockDesign', 'GeometricColorBlock', 'Teccelerates', 'Teccelerates360'),
        )
    with col2:
        image = Image.open('image_ProductTemplate-' + pptstyle + '.png')
        st.image(image, caption=pptstyle)

    product = st.text_input("Enter product name:") 

    if st.button("Generate"):
        ideas = json.loads(insource_idea(product)["result"])
        # st.write(insource_idea(product)["result"])
        # st.write(ideas)
        ppt_gen(pptstyle, product, ideas)

        
if __name__ == "__main__":  
    main() 